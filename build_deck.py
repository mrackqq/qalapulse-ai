"""Генератор pitch-презентации QalaPulse AI (.pptx).

Запуск:  python build_deck.py
Результат: QalaPulse_AI_Pitch.pptx в корне проекта.
Деск брендирован под продукт (indigo / smart city) и опирается на реальные фичи MVP.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- Палитра (совпадает с UI продукта) ---
INK = RGBColor(0x0F, 0x17, 0x2A)
INK2 = RGBColor(0x1E, 0x29, 0x3B)
CANVAS = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
BRAND = RGBColor(0x4F, 0x46, 0xE5)
BRAND_L = RGBColor(0xEE, 0xF2, 0xFF)
BRAND_T = RGBColor(0xA5, 0xB4, 0xFC)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
EMER = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xE1, 0x1D, 0x48)
SLATE_L = RGBColor(0xF1, 0xF5, 0xF9)

HEAD = "Segoe UI Semibold"
HEADB = "Segoe UI Black"
BODY = "Segoe UI"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=CANVAS):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, bg, line=None)
    return s


def rect(s, x, y, w, h, fill, line=None, line_w=1.0, rounded=False, radius=0.06, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def _soft_shadow(shp):
    spPr = shp._element.spPr
    el = spPr.makeelement(qn("a:effectLst"), {})
    sh = el.makeelement(qn("a:outerShdw"), {
        "blurRad": "180000", "dist": "90000", "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "0F172A"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "16000"})
    clr.append(alpha)
    sh.append(clr)
    el.append(sh)
    spPr.append(el)


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    if space_after:
        p.space_after = Pt(space_after)
    for t, ov in runs:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.name = ov.get("font", font)
        r.font.color.rgb = ov.get("color", color)
    return tb


def bullets(s, x, y, w, h, items, size=16, color=INK, gap=10, font=BODY, dot=BRAND):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = "—  "
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.name = font
        r.font.color.rgb = dot
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(size)
        r2.font.name = font
        r2.font.color.rgb = color
    return tb


def eyebrow(s, x, y, label, color=BRAND):
    text(s, x, y, 8, 0.35, label.upper(), size=12.5, color=color, bold=True, font=HEAD)


def page(s, n, title_short):
    text(s, 0.9, 6.92, 8, 0.3, f"QalaPulse AI · {title_short}", size=10, color=MUTED, font=BODY)
    text(s, 11.6, 6.92, 0.9, 0.3, f"{n:02d} / 13", size=10, color=MUTED, font=BODY, align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, label, fill, txt, h=0.42, size=11.5):
    rect(s, x, y, w, h, fill, rounded=True, radius=0.5)
    text(s, x, y, w, h, label, size=size, color=txt, bold=True, font=HEAD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def kpi(s, x, y, w, h, value, label, vcolor=INK, border=LINE, bg=WHITE):
    rect(s, x, y, w, h, bg, line=border, line_w=1.0, rounded=True, radius=0.10, shadow=True)
    text(s, x + 0.22, y + 0.22, w - 0.4, 0.3, label.upper(), size=10.5, color=MUTED, bold=True, font=HEAD)
    text(s, x + 0.2, y + 0.52, w - 0.4, h - 0.6, value, size=30, color=vcolor, bold=True, font=HEADB)


def feature_card(s, x, y, w, h, num, ncolor, nbg, title, body):
    rect(s, x, y, w, h, WHITE, line=LINE, rounded=True, radius=0.07, shadow=True)
    rect(s, x + 0.28, y + 0.3, 0.62, 0.62, nbg, rounded=True, radius=0.3)
    text(s, x + 0.28, y + 0.3, 0.62, 0.62, num, size=17, color=ncolor, bold=True,
         font=HEADB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.3, y + 1.08, w - 0.6, 0.5, title, size=16.5, color=INK, bold=True, font=HEAD)
    text(s, x + 0.3, y + 1.62, w - 0.6, h - 1.7, body, size=12.5, color=MUTED, font=BODY, spacing=1.08)


# ======================================================================
# SLIDE 1 — COVER
# ======================================================================
s = slide(INK)
rect(s, 0, 0, 13.333, 7.5, INK, line=None)
# decorative blobs
rect(s, 8.6, -1.6, 6.5, 6.5, INK2, rounded=True, radius=0.5)
rect(s, 9.7, 4.6, 4.2, 4.2, RGBColor(0x23, 0x1E, 0x55), rounded=True, radius=0.5)
# logo mark
rect(s, 0.9, 0.85, 0.72, 0.72, BRAND, rounded=True, radius=0.28)
text(s, 0.9, 0.85, 0.72, 0.72, "QP", size=20, color=WHITE, bold=True, font=HEADB,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.75, 0.9, 6, 0.4, "QalaPulse AI", size=17, color=WHITE, bold=True, font=HEAD)
text(s, 1.75, 1.28, 6, 0.3, "Astana City Operations", size=11.5, color=BRAND_T, font=BODY)

chip(s, 0.9, 2.55, 4.55, "SMART CITY MVP · АСТАНА", RGBColor(0x26, 0x2C, 0x52), BRAND_T, h=0.46, size=12)
text(s, 0.86, 3.15, 11.2, 1.8,
     "AI-радар городских\nпроблем Астаны", size=52, color=WHITE, bold=True, font=HEADB, spacing=0.98)
text(s, 0.9, 5.15, 9.7, 1.0,
     "От обращения жителя до назначенной службы, SLA и отчёта для акимата. "
     "AI ставит приоритет 0–100, риск и находит дубли — за секунды.",
     size=15.5, color=RGBColor(0xCB, 0xD5, 0xE1), font=BODY, spacing=1.15)

chip(s, 0.9, 6.45, 4.0, "Astana Innovations Accelerator 2026", RGBColor(0x1B, 0x21, 0x3D), RGBColor(0x94, 0xA3, 0xB8), h=0.44, size=11)
text(s, 9.4, 6.5, 3.0, 0.4, "Pitch Deck · v0.2", size=11.5, color=MUTED, font=BODY, align=PP_ALIGN.RIGHT)

# ======================================================================
# SLIDE 2 — ПРОБЛЕМА
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Проблема города")
text(s, 0.86, 1.15, 11.5, 1.2, "Обращений много —\nприоритет не очевиден", size=34, color=INK, bold=True, font=HEADB, spacing=1.0)

cards = [
    ("Поток жалоб", "Дороги, гололёд, подтопления, транспорт, мусор, ЖКХ, экология приходят одновременно и вперемешку."),
    ("Нет приоритета", "Вручную не видно, где риск выше, где жалобы повторяются и какую точку чинить первой."),
    ("Теряется контроль", "Сложно отследить, какая служба отвечает, уложились ли в срок и что реально сделано."),
]
x = 0.9
for i, (t, b) in enumerate(cards):
    rect(s, x, 2.75, 3.74, 2.5, WHITE, line=LINE, rounded=True, radius=0.07, shadow=True)
    rect(s, x, 2.75, 3.74, 0.12, [ROSE, AMBER, BRAND][i], rounded=False)
    text(s, x + 0.32, 3.12, 3.1, 0.5, t, size=17, color=INK, bold=True, font=HEAD)
    text(s, x + 0.32, 3.66, 3.12, 1.5, b, size=13, color=MUTED, font=BODY, spacing=1.1)
    x += 3.96

rect(s, 0.9, 5.6, 11.53, 1.05, INK, rounded=True, radius=0.12)
text(s, 1.3, 5.6, 11, 1.05,
     "«Город генерирует сотни обращений в день. Без приоритизации срочное тонет в потоке рутины.»",
     size=15.5, color=WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
page(s, 2, "Проблема")

# ======================================================================
# SLIDE 3 — РЕШЕНИЕ
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Решение")
text(s, 0.86, 1.15, 11.5, 0.8, "QalaPulse AI — операционный слой города", size=30, color=INK, bold=True, font=HEADB)
text(s, 0.9, 2.0, 11.4, 0.8,
     "Web-платформа, которая превращает текст жалобы в структурированную, приоритизированную задачу "
     "с ответственной службой и сроком SLA.", size=15, color=MUTED, font=BODY, spacing=1.12)

pillars = [
    ("Анализ", "Классификация темы, района, риска и confidence по тексту и геолокации.", BRAND, BRAND_L),
    ("Приоритет", "Единый score 0–100, который выстраивает очередь работ по городу.", ACCENT, RGBColor(0xEC, 0xFE, 0xFF)),
    ("Назначение", "Автовыбор службы и срока SLA, статусы и фото «до/после».", EMER, RGBColor(0xEC, 0xFD, 0xF5)),
    ("Контроль", "Карта, heatmap, рейтинг служб и районов, отчёты для акимата.", AMBER, RGBColor(0xFF, 0xFB, 0xEB)),
]
x = 0.9
for t, b, c, bg in pillars:
    rect(s, x, 3.15, 2.78, 2.9, WHITE, line=LINE, rounded=True, radius=0.08, shadow=True)
    rect(s, x + 0.3, 3.45, 0.7, 0.18, c, rounded=True, radius=0.5)
    text(s, x + 0.3, 3.8, 2.2, 0.5, t, size=18, color=INK, bold=True, font=HEADB)
    text(s, x + 0.3, 4.35, 2.22, 1.6, b, size=12.5, color=MUTED, font=BODY, spacing=1.12)
    x += 2.94
page(s, 3, "Решение")

# ======================================================================
# SLIDE 4 — КАК РАБОТАЕТ
# ======================================================================
s = slide(INK)
eyebrow(s, 0.9, 0.75, "Как работает", BRAND_T)
text(s, 0.86, 1.15, 11.5, 0.8, "От жалобы жителя до закрытой заявки", size=30, color=WHITE, bold=True, font=HEADB)

steps = [
    ("01", "Житель отправляет обращение", "Текст, адрес, точка на карте и фото «до». Всё в web — без приложения и регистрации.", BRAND),
    ("02", "AI ставит приоритет", "Категория, район, риск, priority 0–100, дубли рядом и объяснение балла простым языком.", ACCENT),
    ("03", "Служба берёт в работу", "Автоназначение службы и SLA, смена статусов, фото «после», отчёт для акимата.", EMER),
]
x = 0.9
for n, t, b, c in steps:
    rect(s, x, 2.7, 3.74, 3.3, INK2, rounded=True, radius=0.08)
    rect(s, x + 0.32, 3.0, 0.78, 0.78, c, rounded=True, radius=0.26)
    text(s, x + 0.32, 3.0, 0.78, 0.78, n, size=20, color=WHITE, bold=True, font=HEADB,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.34, 4.0, 3.1, 0.7, t, size=16.5, color=WHITE, bold=True, font=HEAD, spacing=1.0)
    text(s, x + 0.34, 4.75, 3.1, 1.2, b, size=12.5, color=RGBColor(0xCB, 0xD5, 0xE1), font=BODY, spacing=1.12)
    x += 3.96
page(s, 4, "Как работает")

# ======================================================================
# SLIDE 5 — AI-ДВИЖОК
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "AI-движок")
text(s, 0.86, 1.15, 11.5, 0.8, "Прозрачный scoring, который работает локально", size=28, color=INK, bold=True, font=HEADB)

# left: example card
rect(s, 0.9, 2.2, 5.5, 4.2, INK, rounded=True, radius=0.06, shadow=True)
text(s, 1.25, 2.5, 4.9, 0.3, "ПРИМЕР АНАЛИЗА", size=11, color=BRAND_T, bold=True, font=HEAD)
text(s, 1.25, 2.95, 4.9, 1.1,
     "«На переходе у Сарыарка стёрлась разметка и не работает подсветка»",
     size=15.5, color=WHITE, font=BODY, spacing=1.12)
rows = [("Категория", "Безопасность"), ("Priority", "94 / 100  ·  high"),
        ("Служба", "Служба общественной безопасности"), ("SLA", "12 часов")]
yy = 4.25
for k, v in rows:
    text(s, 1.25, yy, 1.7, 0.35, k, size=12.5, color=RGBColor(0x94, 0xA3, 0xB8), font=BODY)
    text(s, 2.95, yy, 3.3, 0.35, v, size=12.5, color=WHITE, bold=True, font=HEAD)
    yy += 0.5

# right: capabilities
caps = [
    "Классификация: 11 городских категорий",
    "Priority score 0–100 + risk level",
    "AI confidence по тексту, гео и фото",
    "Duplicate detection в радиусе 500 м",
    "Объяснение балла (почему срочно)",
    "Rule-based офлайн + опциональный LLM",
]
rect(s, 6.7, 2.2, 5.73, 4.2, WHITE, line=LINE, rounded=True, radius=0.06, shadow=True)
text(s, 7.05, 2.5, 5, 0.4, "Что делает движок", size=17, color=INK, bold=True, font=HEAD)
bullets(s, 7.05, 3.1, 5.1, 3.1, caps, size=13.5, gap=11)
page(s, 5, "AI-движок")

# ======================================================================
# SLIDE 6 — ПРОДУКТ / DASHBOARD
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Продукт")
text(s, 0.86, 1.15, 11.5, 0.8, "Smart city command center", size=30, color=INK, bold=True, font=HEADB)

# mock dashboard frame
rect(s, 0.9, 2.15, 7.2, 4.25, WHITE, line=LINE, rounded=True, radius=0.05, shadow=True)
rect(s, 0.9, 2.15, 7.2, 0.55, SLATE_L, rounded=True, radius=0.05)
for i, c in enumerate([ROSE, AMBER, EMER]):
    rect(s, 1.2 + i * 0.28, 2.36, 0.15, 0.15, c, rounded=True, radius=0.5)
text(s, 2.4, 2.15, 5, 0.55, "qalapulse.ai / dashboard", size=11, color=MUTED, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
# mini kpis
mk = [("84", "Всего", INK), ("29", "High-risk", ROSE), ("7", "Overdue", AMBER)]
mx = 1.2
for v, l, c in mk:
    rect(s, mx, 2.95, 2.1, 0.95, CANVAS, line=LINE, rounded=True, radius=0.12)
    text(s, mx + 0.2, 3.08, 1.8, 0.3, l.upper(), size=9.5, color=MUTED, bold=True, font=HEAD)
    text(s, mx + 0.18, 3.32, 1.8, 0.5, v, size=22, color=c, bold=True, font=HEADB)
    mx += 2.27
# map placeholder
rect(s, 1.2, 4.1, 6.6, 2.05, RGBColor(0xE8, 0xEE, 0xF6), rounded=True, radius=0.05)
for (dx, dy, c) in [(2.0, 4.7, ROSE), (3.1, 5.3, AMBER), (4.3, 4.6, EMER), (5.2, 5.2, ROSE),
                    (6.0, 4.9, AMBER), (3.6, 4.5, EMER), (4.9, 5.5, ROSE)]:
    rect(s, dx, dy, 0.34, 0.34, c, rounded=True, radius=0.5)
text(s, 1.4, 5.78, 5, 0.3, "Маркеры · кластеры · heatmap по 5 районам", size=10.5, color=MUTED, font=BODY)

# right: feature list
feats = [
    "10 KPI: новые, в работе, решено, high-risk, overdue",
    "Карта: маркеры / кластеры / heatmap",
    "Фильтры: район, категория, статус, риск, SLA",
    "Очередь обращений по priority и SLA",
    "Цвет точки = уровень риска",
]
rect(s, 8.35, 2.15, 4.08, 4.25, INK, rounded=True, radius=0.06, shadow=True)
text(s, 8.65, 2.45, 3.5, 0.4, "На одном экране", size=16, color=WHITE, bold=True, font=HEAD)
bullets(s, 8.65, 3.05, 3.5, 3.2, feats, size=12.5, gap=11, color=RGBColor(0xE2, 0xE8, 0xF0), dot=BRAND_T)
page(s, 6, "Продукт")

# ======================================================================
# SLIDE 7 — ДЛЯ КОГО
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Для кого")
text(s, 0.86, 1.15, 11.5, 0.8, "Одна платформа — три аудитории", size=30, color=INK, bold=True, font=HEADB)

aud = [
    ("Жители", "Обращение не теряется", "Понятный номер, статус и объяснение, почему проблема срочная или обычная.", BRAND, BRAND_L),
    ("Службы", "Видно, что делать первым", "Очередь по priority и SLA, кабинет службы, фото «после», закрытие заявок.", ACCENT, RGBColor(0xEC, 0xFE, 0xFF)),
    ("Акимат", "Город в реальном времени", "High-risk зоны, рейтинг служб и районов, статистика и история исполнения.", EMER, RGBColor(0xEC, 0xFD, 0xF5)),
]
x = 0.9
for tag, t, b, c, bg in aud:
    rect(s, x, 2.5, 3.74, 3.6, WHITE, line=LINE, rounded=True, radius=0.07, shadow=True)
    chip(s, x + 0.32, 2.8, 1.7, tag, bg, c, h=0.44, size=12.5)
    text(s, x + 0.32, 3.5, 3.1, 0.9, t, size=18.5, color=INK, bold=True, font=HEADB, spacing=1.0)
    text(s, x + 0.32, 4.55, 3.12, 1.5, b, size=13, color=MUTED, font=BODY, spacing=1.12)
    x += 3.96
page(s, 7, "Для кого")

# ======================================================================
# SLIDE 8 — DEMO СЦЕНАРИЙ
# ======================================================================
s = slide(INK)
eyebrow(s, 0.9, 0.75, "Демо за 90 секунд", BRAND_T)
text(s, 0.86, 1.15, 11.5, 0.8, "Полный loop вживую перед жюри", size=30, color=WHITE, bold=True, font=HEADB)

flow = [
    ("1", "Создать обращение", "Тёмный участок на пешеходном маршруте у остановки."),
    ("2", "Показать AI-решение", "Категория, high risk, priority, служба, SLA, объяснение."),
    ("3", "Найти дубли", "Похожие жалобы рядом — кластер на карте."),
    ("4", "Закрыть loop", "Dashboard → взять в работу → отчёт XLSX/PDF."),
]
x = 0.9
for n, t, b in flow:
    rect(s, x, 2.7, 2.78, 3.0, INK2, rounded=True, radius=0.08)
    rect(s, x + 0.3, 3.0, 0.66, 0.66, BRAND, rounded=True, radius=0.3)
    text(s, x + 0.3, 3.0, 0.66, 0.66, n, size=18, color=WHITE, bold=True, font=HEADB,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.32, 3.85, 2.2, 0.7, t, size=15, color=WHITE, bold=True, font=HEAD, spacing=1.0)
    text(s, x + 0.32, 4.55, 2.2, 1.1, b, size=11.5, color=RGBColor(0xCB, 0xD5, 0xE1), font=BODY, spacing=1.1)
    x += 2.94
text(s, 0.9, 6.05, 11, 0.5, "Кнопка «Загрузить pitch dataset» наполняет систему реалистичными болями Астаны в один клик.",
     size=12.5, color=BRAND_T, font=BODY)
page(s, 8, "Демо")

# ======================================================================
# SLIDE 9 — МЕТРИКИ / ЧТО ВНУТРИ
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Что внутри MVP")
text(s, 0.86, 1.15, 11.5, 0.8, "Реальные цифры демо-датасета", size=30, color=INK, bold=True, font=HEADB)

data = [
    ("84", "обращения в демо", INK),
    ("11", "категорий проблем", BRAND),
    ("5", "районов Астаны", ACCENT),
    ("0–100", "priority score", EMER),
    ("5", "кластеров дублей", AMBER),
    ("24/72/168ч", "SLA по риску", ROSE),
    ("XLSX·CSV·PDF", "форматы отчётов", INK),
    ("100%", "локально, без ключей", BRAND),
]
x, y = 0.9, 2.35
for i, (v, l, c) in enumerate(data):
    kpi(s, x, y, 2.78, 1.75, v, l, vcolor=c)
    x += 2.94
    if i % 4 == 3:
        x = 0.9
        y += 2.0
page(s, 9, "Метрики")

# ======================================================================
# SLIDE 10 — ЦЕННОСТЬ
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Ценность")
text(s, 0.86, 1.15, 11.5, 0.8, "Что меняет QalaPulse AI", size=30, color=INK, bold=True, font=HEADB)
val = [
    ("Срочное не тонет", "Priority и SLA выводят опасные точки наверх автоматически.", BRAND),
    ("Меньше дублей", "Похожие жалобы группируются — служба не делает работу дважды.", ACCENT),
    ("Прозрачность", "Видно ответственного, срок и фактический результат «до/после».", EMER),
    ("Данные для решений", "Рейтинг служб и районов, heatmap и отчёты для акимата.", AMBER),
]
x, y = 0.9, 2.55
for i, (t, b, c) in enumerate(val):
    cx = x + (i % 2) * 5.85
    cy = y + (i // 2) * 1.95
    rect(s, cx, cy, 5.65, 1.7, WHITE, line=LINE, rounded=True, radius=0.08, shadow=True)
    rect(s, cx, cy + 0.28, 0.12, 1.14, c, rounded=True, radius=0.5)
    text(s, cx + 0.42, cy + 0.28, 5.0, 0.5, t, size=17, color=INK, bold=True, font=HEAD)
    text(s, cx + 0.42, cy + 0.82, 5.0, 0.8, b, size=12.5, color=MUTED, font=BODY, spacing=1.1)
page(s, 10, "Ценность")

# ======================================================================
# SLIDE 11 — ТЕХНОЛОГИИ
# ======================================================================
s = slide(INK)
eyebrow(s, 0.9, 0.75, "Технологии", BRAND_T)
text(s, 0.86, 1.15, 11.5, 0.8, "Лёгкий, переносимый стек", size=30, color=WHITE, bold=True, font=HEADB)
stack = [
    ("Backend", "FastAPI · SQLAlchemy · Pydantic"),
    ("Frontend", "Jinja2 · Tailwind CSS · адаптивный UI"),
    ("Карты", "Leaflet · markercluster · heat"),
    ("Аналитика", "Chart.js · KPI · рейтинги"),
    ("Отчёты", "ReportLab PDF · XLSX · CSV · print"),
    ("Данные", "PostgreSQL (prod) · SQLite (локально)"),
    ("AI", "Rule-based движок · опциональный LLM"),
    ("Деплой", "Docker Compose · один запуск"),
]
x, y = 0.9, 2.45
for i, (t, b) in enumerate(stack):
    cx = x + (i % 2) * 5.85
    cy = y + (i // 2) * 1.02
    rect(s, cx, cy, 5.65, 0.86, INK2, rounded=True, radius=0.12)
    text(s, cx + 0.3, cy, 1.9, 0.86, t, size=14, color=BRAND_T, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 2.2, cy, 3.3, 0.86, b, size=12.5, color=WHITE, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
page(s, 11, "Технологии")

# ======================================================================
# SLIDE 12 — ROADMAP
# ======================================================================
s = slide()
eyebrow(s, 0.9, 0.75, "Дальше")
text(s, 0.86, 1.15, 11.5, 0.8, "Дорожная карта", size=30, color=INK, bold=True, font=HEADB)
road = [
    ("Сейчас · MVP", ["Web submit + AI scoring", "Карта, SLA, дубли, отчёты", "Роли: житель / служба / район / акимат"], BRAND),
    ("Next · пилот", ["Интеграция с реестром служб", "Уведомления и эскалации", "Мобильная подача с фото-гео"], ACCENT),
    ("Later · масштаб", ["LLM-анализ и приоритезация", "Открытые данные и API города", "Прогноз горячих зон по сезонам"], EMER),
]
x = 0.9
for t, items, c in road:
    rect(s, x, 2.5, 3.74, 3.7, WHITE, line=LINE, rounded=True, radius=0.07, shadow=True)
    rect(s, x, 2.5, 3.74, 0.13, c)
    text(s, x + 0.32, 2.85, 3.1, 0.5, t, size=16, color=INK, bold=True, font=HEADB)
    bullets(s, x + 0.32, 3.5, 3.15, 2.6, items, size=12.5, gap=10, dot=c)
    x += 3.96
page(s, 12, "Roadmap")

# ======================================================================
# SLIDE 13 — ЗАКРЫТИЕ
# ======================================================================
s = slide(INK)
rect(s, 8.6, -1.8, 7, 7, INK2, rounded=True, radius=0.5)
rect(s, 0.9, 0.9, 0.72, 0.72, BRAND, rounded=True, radius=0.28)
text(s, 0.9, 0.9, 0.72, 0.72, "QP", size=20, color=WHITE, bold=True, font=HEADB,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.75, 0.95, 6, 0.5, "QalaPulse AI", size=17, color=WHITE, bold=True, font=HEAD)

text(s, 0.86, 2.7, 11.5, 1.7, "Готовый smart city MVP,\nкоторый не стыдно показать", size=40, color=WHITE, bold=True, font=HEADB, spacing=1.0)
text(s, 0.9, 4.55, 9.5, 1.0,
     "От жалобы жителя до отчёта акимата — один продукт, один loop. "
     "Запускается одной командой, работает локально.", size=15.5, color=RGBColor(0xCB, 0xD5, 0xE1), font=BODY, spacing=1.15)

chip(s, 0.9, 5.85, 3.0, "docker compose up", INK2, BRAND_T, h=0.5, size=12.5)
chip(s, 4.05, 5.85, 2.3, "/demo · /dashboard", INK2, BRAND_T, h=0.5, size=12.5)
text(s, 0.9, 6.7, 11, 0.4, "Astana Innovations Accelerator 2026  ·  Спасибо за внимание",
     size=12, color=MUTED, font=BODY)

prs.save("QalaPulse_AI_Pitch.pptx")
print(f"OK · {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> QalaPulse_AI_Pitch.pptx")
